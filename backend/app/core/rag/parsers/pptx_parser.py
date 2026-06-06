"""PPTX parser — extracts text from PowerPoint slides."""
from app.core.rag.parsers.base import BaseParser


class PptxParser(BaseParser):
    supported_extensions = [".pptx", ".ppt"]

    def parse(self, file_path: str) -> str:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            lines.append(para.text.strip())
            if lines:
                slides.append(f"Slide {i + 1}:\n" + "\n".join(lines))
        return "\n\n".join(slides)
